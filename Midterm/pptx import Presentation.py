from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize Presentation
prs = Presentation()

# Data for the slides based on your script
slides_data = [
    {"title": "AI Deepfakes in Entertainment", "content": "Hello everyone. Today I will be presenting a case study analysis..."},
    {"title": "The Problem Deepfake AI Addresses", "content": "Film studios constantly face challenges such as aging actors..."},
    {"title": "How Deepfake AI Works", "content": "Deepfake systems are powered mainly by deep learning models..."},
    {"title": "Use in Film Production", "content": "Major film studios have already adopted AI-assisted face technologies..."},
    {"title": "Benefits of Deepfake AI", "content": "There are several major advantages. First, creative flexibility..."},
    {"title": "Technical and Operational Challenges", "content": "Despite its promise, deepfake AI still faces important limitations..."},
    {"title": "Ethical Issues", "content": "Deepfake AI raises several serious ethical concerns..."},
    {"title": "Broader Societal Effects", "content": "On the positive side, deepfake AI enables new storytelling possibilities..."},
    {"title": "Where the Technology Is Going", "content": "Looking ahead, we expect several important developments..."},
    {"title": "Recommendations and Final Thoughts", "content": "To ensure responsible use of deepfake AI, the industry should..."}
]

for data in slides_data:
    # Add a blank slide layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.add_paragraph()
    p.text = data["title"]
    p.font.bold = True
    p.font.size = Pt(32)

    # Placeholder for Image (Left Side)
    # Note: Replace 'placeholder.jpg' with your actual image paths
    try:
        slide.shapes.add_picture('placeholder.jpg', Inches(0.5), Inches(1.5), width=Inches(4))
    except:
        # If no image found, create a rectangle as a placeholder
        slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(4)).text = "[Image Placeholder]"

    # Add Text (Right Side)
    content_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    p = content_frame.add_paragraph()
    p.text = data["content"]
    p.font.size = Pt(18)

# Save the presentation
prs.save('AI_Deepfakes_Presentation.pptx')
print("Presentation created successfully!")
